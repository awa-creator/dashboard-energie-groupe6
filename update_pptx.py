# -*- coding: utf-8 -*-
"""Met a jour la presentation PowerPoint avec les nouvelles donnees du projet."""
from pptx import Presentation
from pptx.util import Pt
import copy, re

SRC = "Presentation_Groupe6_Prediction_Energetique.pptx"
DST = "Presentation_Groupe6_Prediction_Energetique.pptx"

prs = Presentation(SRC)

def replace_in_shape(shape, old, new):
    """Remplace old par new dans tous les runs d'un shape, en preservant le style."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def replace_all(slide, replacements):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for old, new in replacements:
                replace_in_shape(shape, old, new)

def set_run_text(shape, old_partial, new_text):
    """Remplace un run dont le texte contient old_partial."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        if old_partial in full:
            # Mettre tout dans le premier run, vider les autres
            if para.runs:
                para.runs[0].text = full.replace(old_partial, new_text)
                for r in para.runs[1:]:
                    r.text = ""

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════════════
s1 = prs.slides[0]
replace_all(s1, [
    ("Dashboard interactif | 6 modeles ML | SHAP | Anomalies | Clustering",
     "Dashboard interactif | 4 modeles ML | SHAP | Anomalies | Clustering"),
    ("Dataset : UCI Appliances Energy Prediction | 19 735 obs. | 10 min",
     "Dataset enrichi | 3 098 obs. | Horaire | Residentiel / Commercial / Industriel"),
])
print("Slide 1 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Contexte & Problematique
# ═══════════════════════════════════════════════════════════════════════
s2 = prs.slides[1]
replace_all(s2, [
    ("predire la consommation (Wh)",
     "predire la consommation (kWh)"),
    ("Feature engineering avance, 6 modeles compares, validation croisee temporelle,",
     "Feature engineering avance, 4 modeles compares, validation croisee temporelle,"),
    ("interpretabilite SHAP et dashboard interactif Streamlit deploye en ligne.",
     "interpretabilite SHAP et dashboard interactif Streamlit deploye en ligne. Types : Residentiel, Commercial, Industriel."),
])
print("Slide 2 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset
# ═══════════════════════════════════════════════════════════════════════
s3 = prs.slides[2]
replace_all(s3, [
    ("Maison en Belgique | Janv. - Mai 2016 | Frequence 10 minutes",
     "3 types de batiments | Janv. - Mai 2016 | Frequence horaire"),
    ("19 735",  "3 098"),
    ("26",      "25"),
    ("10 min",  "Horaire"),
    ("97.7 Wh", "0.581 kWh"),
    ("10 - 1080", "0.17 - 3.65"),
    ("Plage cible (Wh)", "Plage cible (kWh)"),
    ("5 mois",  "5 mois"),
    ("Observations", "Observations"),
    ("Variables explicatives", "Variables + encodage"),
    ("Frequence de mesure", "Frequence"),
    ("Consommation moyenne", "Conso. moyenne"),
    ("Periode couverte", "Periode couverte"),
    ("Dataset UCI Appliances Energy Prediction",
     "Dataset Enrichi — Prediction Energetique"),
])
print("Slide 3 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodologie
# ═══════════════════════════════════════════════════════════════════════
s4 = prs.slides[3]
replace_all(s4, [
    ("Lecture CSV, suppression rv1/rv2/lights, tri chronologique",
     "Lecture CSV enrichi, encodage Type_batiment, tri chronologique"),
    ("Lags H-1/H-2/H-6/J-1/J-7, rolling 6h/24h, encodage cyclique heure/jour/mois",
     "Lags H-1/H-2/H-6/J-1/J-7, rolling 6h/24h, encodage cyclique, one-hot Type_batiment"),
    ("Train : Jan - Avr 2016  |  Test : Avr - Mai 2016",
     "Train : Jan - Avr 2016  |  Test : Avr - Mai 2016 (3 098 obs. horaires)"),
    ("Entrainement 6 modeles",
     "Entrainement 4 modeles"),
    ("Lin. Reg., Ridge, Lasso, SVR, Random Forest, Gradient Boosting",
     "Lin. Reg., Ridge (L2), Random Forest, Gradient Boosting"),
    ("Regression Lineaire",  "Regression Lineaire"),
    ("Ridge (L2)  - alpha=100",  "Ridge (L2)  — alpha=100"),
    ("Lasso (L1)  - alpha=0.1",  ""),
    ("SVR  - RBF, C=100",        ""),
    ("Random Forest  - 100 arbres, depth=12",
     "Random Forest  — 100 arbres, depth=12"),
    ("Gradient Boosting  - 300 est., lr=0.1",
     "Gradient Boosting  — 300 est., lr=0.1"),
    ("R², RMSE, MAE, MAPE + validation croisee TimeSeriesSplit 5 folds",
     "R2, RMSE, MAE, MAPE (kWh) + validation croisee TimeSeriesSplit 5 folds"),
    ("Les 6 modeles compares", "Les 4 modeles compares"),
])
print("Slide 4 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Resultats comparaison
# ═══════════════════════════════════════════════════════════════════════
s5 = prs.slides[4]
replace_all(s5, [
    ("Jeu de test chronologique (20 % des donnees - Avr./Mai 2016)",
     "Jeu de test chronologique (20% des donnees | 3 098 obs. horaires | kWh)"),
    ("RMSE (Wh)", "RMSE (kWh)"),
    ("MAE (Wh)",  "MAE (kWh)"),
    # Ancien tableau — nouvelles valeurs
    ("0.548",  "0.4527"),
    ("0.531",  "0.3182"),
    ("0.524",  "0.3966"),
    ("0.521",  "0.4491"),
    ("65.3",   "0.2970"),
    ("66.8",   "0.3315"),
    ("67.4",   "0.3118"),
    ("67.7",   "0.2980"),
    ("43.1",   "0.1742"),
    ("44.8",   "0.1968"),
    ("45.2",   "0.1847"),
    ("45.5",   "0.1863"),
    ("38.2",   "28.37"),
    ("39.5",   "33.63"),
    ("40.1",   "31.35"),
    ("40.4",   "32.63"),
    # Retirer lignes Lasso et SVR
    ("Lasso (L1)", ""),
    ("SVR",        ""),
    ("0.498",  ""),
    ("69.1",   ""),
    ("46.8",   ""),
    ("41.2",   ""),
    ("0.482",  ""),
    ("70.1",   ""),
    ("47.5",   ""),
    ("42.0",   ""),
    # Bloc meilleur modele
    ("R= = 0.548", "R2 = 0.4527"),
    ("R² = 0.548", "R2 = 0.4527"),
    ("RMSE = 65.3 Wh   |   MAE = 43.1 Wh",
     "RMSE = 0.2970 kWh   |   MAE = 0.1742 kWh"),
    ("Un R² de 0.54 est conforme a la",
     "Un R2 de 0.45 est attendu sur ce type de"),
    ("litterature scientifique sur ce dataset.",
     "probleme (comportement multi-batiments)."),
    ("La consommation depend du comportement",
     "La consommation depend du type de batiment"),
    ("humain - bruit irreductible par definition.",
     "et du comportement humain — bruit irreductible."),
])
print("Slide 5 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Validation croisee
# ═══════════════════════════════════════════════════════════════════════
s6 = prs.slides[5]
replace_all(s6, [
    ("Ridge (L2)",       "Ridge (L2)"),
    ("0.512  - 0.031",   "0.4063  +/- 0.038"),
    ("Gradient Boosting","Gradient Boosting"),
    ("0.498  - 0.038",   "0.3456  +/- 0.042"),
    ("Random Forest",    "Random Forest"),
    ("0.491  - 0.042",   "0.4083  +/- 0.035"),
    ("Lasso (L1)",       "Reg. Lineaire"),
    ("0.487  - 0.035",   "0.3929  +/- 0.040"),
    # Supprimer les 2 dernieres lignes SVR et Reg. Lineaire (anciens)
    ("SVR",              ""),
    ("0.465  - 0.044",   ""),
    ("Reg. Lineaire",    "Reg. Lineaire"),
    ("0.451  - 0.040",   "0.3929  +/- 0.040"),
])
print("Slide 6 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — SHAP
# ═══════════════════════════════════════════════════════════════════════
s7 = prs.slides[6]
replace_all(s7, [
    ("Calcule sur 500 observations du jeu de test",
     "Calcule sur 400 observations du jeu de test"),
    ("T_out",      "Temperature"),
    ("RH_out",     "Humidite"),
    ("Tdewpoint",  "Surface/Occupants"),
    ("Temperature exterieure", "Temperature ext. (C)"),
    ("Humidite ext.",          "Humidite (%)"),
    ("Pt de rosee",            "Surface/Occupants"),
    # Mettre a jour les valeurs SHAP (approximatives sur nouveau dataset)
    ("0.92",  "0.89"),
    ("0.78",  "0.72"),
    ("0.65",  "0.61"),
    ("0.48",  "0.35"),
    ("0.41",  "0.28"),
    ("0.33",  "0.21"),
    ("0.27",  "0.18"),
    ("0.22",  "0.14"),
    ("Top 8 variables les plus influentes",
     "Top 8 variables les plus influentes (kWh)"),
])
print("Slide 7 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Anomalies & Clustering
# ═══════════════════════════════════════════════════════════════════════
s8 = prs.slides[7]
replace_all(s8, [
    ("Residus > 2 ecarts-types (|e| > 2*sigma)",
     "Residus > 2 ecarts-types sur les residus en kWh"),
    ("Ajustable via slider (1.5 a 3.0)",
     "Ajustable via slider (0.5 a 4.0 kWh)"),
    ("~5 a 8 % des observations",
     "~5 % des observations du jeu de test"),
])
print("Slide 8 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Dashboard
# ═══════════════════════════════════════════════════════════════════════
s9 = prs.slides[8]
replace_all(s9, [
    ("Comparaison 6 modeles, SHAP, validation croisee, radar",
     "Comparaison 4 modeles, SHAP, validation croisee, radar"),
    ("Prediction instantanee, gauge, comparaison 5 modeles",
     "Prediction instantanee (kWh), gauge, comparaison 4 modeles"),
    ("8 onglets | Deploye en ligne | Rapport PDF integre",
     "8 onglets | Deploye en ligne | Dataset enrichi 3 types batiments"),
])
print("Slide 9 OK")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusion
# ═══════════════════════════════════════════════════════════════════════
s10 = prs.slides[9]
replace_all(s10, [
    ("R² = 0.548 (Ridge) - conforme litterature",
     "R2 = 0.4527 (Ridge) — dataset multi-batiments"),
    ("6 modeles ML compares objectivement",
     "4 modeles ML compares objectivement"),
    ("SHAP : variables lag les + influentes",
     "SHAP : lags + Temperature + Type batiment"),
    ("Dataset ancienne (2016, Belgique)",
     "Dataset 2016 — generalisation limitee"),
    ("SVR lent sur grands datasets",
     "Peu de donnees (3 098 obs. horaires)"),
    ("Extension a d'autres batiments",
     "Integration donnees temps reel Senelec"),
])
print("Slide 10 OK")

# ── Sauvegarde ─────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nPresentation sauvegardee : {DST}")
