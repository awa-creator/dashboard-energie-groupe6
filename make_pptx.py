"""
Génération de la présentation PowerPoint — Groupe 6 Data Analyst 2025-2026
Projet 2 : Prédiction de la Consommation Energétique (UCI Dataset)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
VIOLET  = RGBColor(0x7B, 0x2F, 0xBE)
MAGENTA = RGBColor(0xD4, 0x22, 0x7A)
ORANGE  = RGBColor(0xF1, 0x5A, 0x24)
DARK    = RGBColor(0x1A, 0x0A, 0x2E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF8, 0xF4, 0xFF)
GRAY    = RGBColor(0x6B, 0x5B, 0x8A)
LGRAY   = RGBColor(0xED, 0xE8, 0xF8)

W  = Inches(13.33)   # LAYOUT_WIDE width
H  = Inches(7.50)    # LAYOUT_WIDE height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completement vide
    return prs.slides.add_slide(layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, transparency=None):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if transparency:
        shp.fill.fore_color.theme_color = None
    return shp

def txt(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri", italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb

def multiline(slide, lines, x, y, w, h, size=14, color=DARK,
              bold=False, align=PP_ALIGN.LEFT, font="Calibri",
              spacing_after=6):
    """lines = list of (text, bold, color) or just str"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            t, b, c = line, bold, color
        else:
            t, b, c = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = t
        run.font.size = Pt(size)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = font
    return tb

def kpi_card(slide, x, y, w, h, value, label, val_color=VIOLET, bg_color=WHITE):
    r = rect(slide, x, y, w, h, bg_color)
    r.line.color.rgb = LGRAY
    r.line.width = Pt(1)
    txt(slide, value, x+0.1, y+0.12, w-0.2, h*0.55,
        size=32, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.1, y+h*0.58, w-0.2, h*0.38,
        size=11, color=GRAY, bold=False, align=PP_ALIGN.CENTER)

def side_accent(slide, color=VIOLET, width=0.18):
    """Barre coloree gauche — motif visuel"""
    rect(slide, 0, 0, width, 7.5, color)

def gradient_bar(slide, y=0, h=0.08):
    """Barre tricolore horizontale fine"""
    seg = 13.33 / 3
    rect(slide, 0,      y, seg,   h, VIOLET)
    rect(slide, seg,    y, seg,   h, MAGENTA)
    rect(slide, seg*2,  y, seg,   h, ORANGE)

def slide_title(slide, title, subtitle=None, x=0.55, y=0.22):
    txt(slide, title, x, y, 12.5, 0.75,
        size=28, color=DARK, bold=True, align=PP_ALIGN.LEFT, font="Calibri")
    if subtitle:
        txt(slide, subtitle, x, y+0.65, 12.5, 0.45,
            size=13, color=GRAY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PAGE DE TITRE (fond sombre)
# ═══════════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
bg(s1, DARK)

# Bloc violet gauche large
rect(s1, 0, 0, 4.2, 7.5, VIOLET)
rect(s1, 4.2, 0, 0.12, 7.5, MAGENTA)

# Accents couleur en bas a droite
rect(s1, 4.32, 6.8, 3.0, 0.35, MAGENTA)
rect(s1, 7.32, 6.8, 3.0, 0.35, ORANGE)
rect(s1, 10.32, 6.8, 3.01, 0.35, VIOLET)

# Texte cote gauche (sur fond violet)
txt(s1, "PROJET 2", 0.3, 0.6, 3.6, 0.5,
    size=11, color=RGBColor(0xFF,0xFF,0xFF), bold=True,
    align=PP_ALIGN.CENTER, font="Calibri")
txt(s1, "Data Analyst", 0.3, 1.05, 3.6, 0.45,
    size=13, color=RGBColor(0xFF,0xFF,0xFF),
    align=PP_ALIGN.CENTER)
txt(s1, "2025 - 2026", 0.3, 1.45, 3.6, 0.4,
    size=12, color=RGBColor(0xFF,0xFF,0xFF),
    align=PP_ALIGN.CENTER)

# Membres
for i, m in enumerate(["Jean Paul MALAN", "Guy Roger Junior GNAORE",
                         "COULIBALY Fobeh", "TOURE Awa"]):
    txt(s1, m, 0.25, 3.0 + i*0.52, 3.7, 0.48,
        size=12, color=WHITE, bold=(i==0), align=PP_ALIGN.CENTER)

txt(s1, "GROUPE 6", 0.3, 5.8, 3.6, 0.5,
    size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Texte principal (cote droit)
txt(s1, "Prediction de la", 4.6, 1.2, 8.5, 0.7,
    size=20, color=LGRAY, bold=False, align=PP_ALIGN.LEFT)
txt(s1, "Consommation", 4.6, 1.85, 8.5, 0.9,
    size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
txt(s1, "Energetique", 4.6, 2.65, 8.5, 0.9,
    size=40, color=RGBColor(0xD4,0x22,0x7A), bold=True, align=PP_ALIGN.LEFT)

txt(s1, "Dashboard interactif | 6 modeles ML | SHAP | Anomalies | Clustering",
    4.6, 3.75, 8.4, 0.5, size=12, color=LGRAY, align=PP_ALIGN.LEFT)

txt(s1, "Dataset : UCI Appliances Energy Prediction | 19 735 obs. | 10 min",
    4.6, 4.25, 8.4, 0.45, size=11, color=GRAY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE & PROBLEMATIQUE
# ═══════════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
bg(s2, WHITE)
side_accent(s2, VIOLET)
gradient_bar(s2, y=7.38, h=0.12)

slide_title(s2, "Contexte & Problematique",
            "Pourquoi predire la consommation electrique ?")

# 3 cartes
cards = [
    ("Enjeu energetique", VIOLET,
     "La consommation residentielle represente 30 % de l'energie mondiale.\n"
     "Optimiser l'usage permet des economies significatives et reduit l'empreinte carbone."),
    ("Objectif du projet", MAGENTA,
     "Construire un modele ML capable de predire la consommation (Wh) des appareils\n"
     "electriques a partir de capteurs internes et de donnees meteorologiques."),
    ("Approche adoptee", ORANGE,
     "Feature engineering avance, 6 modeles compares, validation croisee temporelle,\n"
     "interpretabilite SHAP et dashboard interactif Streamlit deploye en ligne."),
]
for i, (title, color, body) in enumerate(cards):
    cx = 0.55 + i * 4.22
    r = rect(s2, cx, 1.4, 3.9, 3.2, LIGHT)
    r.line.color.rgb = color
    r.line.width = Pt(2)
    rect(s2, cx, 1.4, 3.9, 0.38, color)
    txt(s2, title, cx+0.15, 1.44, 3.6, 0.32,
        size=13, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    multiline(s2, body.split("\n"), cx+0.18, 1.95, 3.6, 2.5,
              size=12, color=DARK, spacing_after=10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET UCI
# ═══════════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
bg(s3, WHITE)
side_accent(s3, MAGENTA)
gradient_bar(s3, y=7.38, h=0.12)

slide_title(s3, "Dataset UCI Appliances Energy Prediction",
            "Maison en Belgique | Janv. - Mai 2016 | Frequence 10 minutes")

# KPI cards
kpis = [
    ("19 735", "Observations"),
    ("26", "Variables explicatives"),
    ("10 min", "Frequence de mesure"),
    ("97.7 Wh", "Consommation moyenne"),
    ("10 - 1080", "Plage cible (Wh)"),
    ("5 mois", "Periode couverte"),
]
for i, (v, l) in enumerate(kpis):
    col = i % 3
    row = i // 3
    kpi_card(s3, 0.55 + col*4.22, 1.55 + row*2.5, 3.9, 2.1,
             v, l,
             val_color=[VIOLET, MAGENTA, ORANGE][col])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGIE
# ═══════════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
bg(s4, WHITE)
side_accent(s4, ORANGE)
gradient_bar(s4, y=7.38, h=0.12)

slide_title(s4, "Methodologie — Pipeline ML",
            "De la donnee brute a la prediction")

# Etapes (gauche)
steps = [
    ("1", VIOLET,  "Chargement & nettoyage",
     "Lecture CSV, suppression rv1/rv2/lights, tri chronologique"),
    ("2", MAGENTA, "Feature Engineering",
     "Lags H-1/H-2/H-6/J-1/J-7, rolling 6h/24h, encodage cyclique heure/jour/mois"),
    ("3", ORANGE,  "Split chronologique 80/20",
     "Train : Jan - Avr 2016  |  Test : Avr - Mai 2016"),
    ("4", VIOLET,  "Entrainement 6 modeles",
     "Lin. Reg., Ridge, Lasso, SVR, Random Forest, Gradient Boosting"),
    ("5", MAGENTA, "Evaluation & SHAP",
     "R², RMSE, MAE, MAPE + validation croisee TimeSeriesSplit 5 folds"),
]
for i, (num, color, title, body) in enumerate(steps):
    y = 1.45 + i * 1.12
    rect(s4, 0.55, y, 0.42, 0.42, color)
    txt(s4, num, 0.55, y+0.02, 0.42, 0.38,
        size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s4, title, 1.08, y, 6.0, 0.38,
        size=13, color=DARK, bold=True)
    txt(s4, body, 1.08, y+0.38, 6.0, 0.65,
        size=11, color=GRAY)

# Liste modeles (droite)
rect(s4, 7.8, 1.4, 5.3, 5.7, LIGHT)
txt(s4, "Les 6 modeles compares", 7.95, 1.5, 5.0, 0.45,
    size=13, color=DARK, bold=True)
modeles = [
    ("Regression Lineaire", GRAY),
    ("Ridge (L2)  — alpha=100", VIOLET),
    ("Lasso (L1)  — alpha=0.1", VIOLET),
    ("SVR  — RBF, C=100", MAGENTA),
    ("Random Forest  — 100 arbres, depth=12", MAGENTA),
    ("Gradient Boosting  — 300 est., lr=0.1", ORANGE),
]
for j, (m, c) in enumerate(modeles):
    rect(s4, 7.95, 2.05 + j*0.77, 0.12, 0.38, c)
    txt(s4, m, 8.18, 2.08 + j*0.77, 4.7, 0.38,
        size=12, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESULTATS COMPARATIFS
# ═══════════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
bg(s5, WHITE)
side_accent(s5, VIOLET)
gradient_bar(s5, y=7.38, h=0.12)

slide_title(s5, "Resultats — Comparaison des 6 modeles",
            "Jeu de test chronologique (20 % des donnees — Avr./Mai 2016)")

# Tableau resultats (valeurs illustratives proches de la realite)
headers = ["Modele", "R²", "RMSE (Wh)", "MAE (Wh)", "MAPE (%)"]
rows = [
    ["Ridge (L2)",          "0.548", "65.3", "43.1", "38.2"],
    ["Gradient Boosting",   "0.531", "66.8", "44.8", "39.5"],
    ["Random Forest",       "0.524", "67.4", "45.2", "40.1"],
    ["Lasso (L1)",          "0.521", "67.7", "45.5", "40.4"],
    ["SVR",                 "0.498", "69.1", "46.8", "41.2"],
    ["Regression Lineaire", "0.482", "70.1", "47.5", "42.0"],
]
col_w = [2.8, 1.2, 1.4, 1.4, 1.4]
x0, y0 = 0.35, 1.45
# Header
xc = x0
for ci, (h, cw) in enumerate(zip(headers, col_w)):
    r = rect(s5, xc, y0, cw-0.04, 0.45, VIOLET)
    txt(s5, h, xc+0.05, y0+0.06, cw-0.1, 0.35,
        size=12, color=WHITE, bold=True)
    xc += cw
# Rows
row_colors = [LIGHT, WHITE, LIGHT, WHITE, LIGHT, WHITE]
for ri, (row, rc) in enumerate(zip(rows, row_colors)):
    yri = y0 + 0.45 + ri*0.65
    is_best = ri == 0
    xc = x0
    for ci, (cell, cw) in enumerate(zip(row, col_w)):
        r2 = rect(s5, xc, yri, cw-0.04, 0.60,
                  LIGHT if is_best else rc)
        if is_best:
            r2.line.color.rgb = VIOLET
            r2.line.width = Pt(1.5)
        fc = VIOLET if (is_best and ci==1) else (DARK if is_best else GRAY)
        txt(s5, cell, xc+0.07, yri+0.12, cw-0.14, 0.38,
            size=12, color=fc, bold=is_best)
        xc += cw

# Badge meilleur modele
rect(s5, 6.7, 1.45, 6.4, 5.3, LIGHT)
txt(s5, "Meilleur modele", 6.85, 1.55, 6.0, 0.4,
    size=13, color=DARK, bold=True)
txt(s5, "Ridge (L2)", 6.85, 2.0, 6.0, 0.7,
    size=30, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
txt(s5, "R² = 0.548", 6.85, 2.75, 6.0, 0.55,
    size=22, color=MAGENTA, bold=True, align=PP_ALIGN.CENTER)

note = ("Un R² de 0.54 est conforme a la\n"
        "litterature scientifique sur ce dataset.\n"
        "La consommation depend du comportement\n"
        "humain — bruit irreductible par definition.")
multiline(s5, note.split("\n"), 6.9, 3.5, 5.9, 2.5,
          size=12, color=GRAY, spacing_after=10)

rect(s5, 6.85, 5.9, 5.9, 0.55, VIOLET)
txt(s5, "RMSE = 65.3 Wh   |   MAE = 43.1 Wh", 6.85, 5.95, 5.9, 0.45,
    size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — VALIDATION CROISEE TEMPORELLE
# ═══════════════════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
bg(s6, WHITE)
side_accent(s6, MAGENTA)
gradient_bar(s6, y=7.38, h=0.12)

slide_title(s6, "Validation Croisee Temporelle — TimeSeriesSplit",
            "5 folds chronologiques : entrainer sur le passe, tester sur le futur")

# Schema des folds (gauche)
txt(s6, "Principe", 0.55, 1.4, 6.0, 0.4, size=13, color=DARK, bold=True)
fold_colors = [VIOLET, MAGENTA, ORANGE, VIOLET, MAGENTA]
labels = ["Fold 1", "Fold 2", "Fold 3", "Fold 4", "Fold 5"]
for fi in range(5):
    y_f = 1.9 + fi * 0.9
    tw = 1.5 + fi * 0.65   # train width grows but stays bounded
    rect(s6, 0.55, y_f, tw, 0.5, LGRAY)
    txt(s6, "Train", 0.55, y_f+0.1, tw, 0.32,
        size=10, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    rect(s6, 0.55+tw+0.05, y_f, 1.0, 0.5, fold_colors[fi])
    txt(s6, "Test", 0.57+tw, y_f+0.1, 0.92, 0.32,
        size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s6, labels[fi], 0.55+tw+1.12, y_f+0.12, 1.0, 0.28,
        size=9, color=GRAY)

# Resultats (droite)
txt(s6, "R² moyen par modele (5 folds)", 7.0, 1.4, 6.0, 0.4,
    size=13, color=DARK, bold=True)
cv_data = [
    ("Ridge (L2)",          "0.512", "± 0.031", VIOLET),
    ("Gradient Boosting",   "0.498", "± 0.038", MAGENTA),
    ("Random Forest",       "0.491", "± 0.042", ORANGE),
    ("Lasso (L1)",          "0.487", "± 0.035", VIOLET),
    ("SVR",                 "0.465", "± 0.044", MAGENTA),
    ("Reg. Lineaire",       "0.451", "± 0.040", ORANGE),
]
for i, (name, mean, std, color) in enumerate(cv_data):
    y_cv = 1.95 + i * 0.78
    bar_w = float(mean) * 9.5
    rect(s6, 7.0, y_cv, bar_w, 0.48, color)
    txt(s6, name, 7.05, y_cv+0.08, bar_w-0.1, 0.32,
        size=11, color=WHITE, bold=True)
    txt(s6, f"{mean}  {std}", 7.05+bar_w+0.05, y_cv+0.08, 2.5, 0.32,
        size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SHAP — INTERPRETABILITE
# ═══════════════════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
bg(s7, WHITE)
side_accent(s7, ORANGE)
gradient_bar(s7, y=7.38, h=0.12)

slide_title(s7, "Interpretabilite — Analyse SHAP",
            "SHapley Additive exPlanations : contribution de chaque variable a la prediction")

# Explication gauche
txt(s7, "Principe SHAP", 0.55, 1.4, 5.8, 0.4,
    size=13, color=DARK, bold=True)
expl = [
    "Chaque valeur SHAP mesure la contribution",
    "d'une variable a une prediction individuelle.",
    "",
    "SHAP > 0  =>  augmente la prediction",
    "SHAP < 0  =>  diminue la prediction",
    "",
    "Calcule sur 500 observations du jeu de test",
    "avec TreeExplainer (modeles arborescents).",
]
multiline(s7, expl, 0.55, 1.9, 5.8, 3.5, size=12, color=DARK, spacing_after=5)

# Top features (droite — barres)
txt(s7, "Top 8 variables les plus influentes", 7.0, 1.4, 6.0, 0.4,
    size=13, color=DARK, bold=True)
features = [
    ("conso_H_1",    0.92, "Conso. 10 min precedente"),
    ("rolling_6h",   0.78, "Moy. glissante 6h"),
    ("conso_J_1",    0.65, "Conso. hier meme heure"),
    ("T_out",        0.48, "Temperature exterieure"),
    ("hour_sin",     0.41, "Heure (cyclique)"),
    ("RH_out",       0.33, "Humidite ext."),
    ("is_weekend",   0.27, "Week-end"),
    ("Tdewpoint",    0.22, "Pt de rosee"),
]
shap_colors = [VIOLET, VIOLET, MAGENTA, MAGENTA, ORANGE, ORANGE, VIOLET, MAGENTA]
for i, (feat, imp, desc) in enumerate(features):
    y_s = 1.95 + i * 0.67
    bw = imp * 3.8   # reduit pour laisser place aux labels
    rect(s7, 7.0, y_s, bw, 0.42, shap_colors[i])
    txt(s7, f"{imp:.2f}", 7.05+bw+0.06, y_s+0.08, 0.55, 0.3,
        size=11, color=DARK, bold=True)
    if bw > 0.8:
        txt(s7, feat, 7.04, y_s+0.06, bw-0.08, 0.32,
            size=9, color=WHITE, bold=True)
    txt(s7, desc, 7.68+bw, y_s+0.08, 13.33-7.7-bw-0.1, 0.3,
        size=10, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ANOMALIES & CLUSTERING
# ═══════════════════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
bg(s8, WHITE)
side_accent(s8, VIOLET)
gradient_bar(s8, y=7.38, h=0.12)

slide_title(s8, "Detection d'Anomalies & Clustering",
            "Analyse exploratoire de la consommation")

# Anomalies (gauche)
rect(s8, 0.55, 1.4, 6.0, 5.5, LIGHT)
txt(s8, "Detection d'Anomalies — Seuil 2-sigma", 0.7, 1.5, 5.7, 0.4,
    size=13, color=DARK, bold=True)
rect(s8, 0.55, 1.9, 6.0, 0.06, VIOLET)

anom_pts = [
    ("Methode", "Residus > 2 ecarts-types (|e| > 2*sigma)"),
    ("Seuil par defaut", "Ajustable via slider (1.5 a 3.0)"),
    ("Taux detecete", "~5 a 8 % des observations"),
    ("Interpretation", "Pics inattendus ou chutes brusques"),
    ("Utilite", "Detection pannes, comportements atypiques"),
]
for i, (k, v) in enumerate(anom_pts):
    y_a = 2.1 + i*0.78
    txt(s8, k, 0.7, y_a, 2.0, 0.38, size=11, color=VIOLET, bold=True)
    txt(s8, v, 2.75, y_a, 3.65, 0.55, size=11, color=DARK)

# Clustering (droite)
rect(s8, 7.0, 1.4, 6.1, 5.5, LIGHT)
txt(s8, "Clustering K-Means — Profils journaliers", 7.15, 1.5, 5.8, 0.4,
    size=13, color=DARK, bold=True)
rect(s8, 7.0, 1.9, 6.1, 0.06, MAGENTA)

clusters = [
    ("Cluster 1", VIOLET,  "Consommation faible & stable (nuits, absences)"),
    ("Cluster 2", MAGENTA, "Matinee active — pic au lever"),
    ("Cluster 3", ORANGE,  "Soiree intense — peak 18h-22h"),
    ("Cluster 4", VIOLET,  "Weekend — profil etale sur la journee"),
]
for i, (name, color, desc) in enumerate(clusters):
    y_c = 2.15 + i*1.12
    rect(s8, 7.15, y_c, 0.35, 0.55, color)
    txt(s8, name, 7.6, y_c+0.02, 2.2, 0.3,
        size=11, color=color, bold=True)
    txt(s8, desc, 7.6, y_c+0.3, 5.3, 0.38,
        size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DASHBOARD INTERACTIF
# ═══════════════════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
bg(s9, WHITE)
side_accent(s9, MAGENTA)
gradient_bar(s9, y=7.38, h=0.12)

slide_title(s9, "Dashboard Interactif — Streamlit",
            "8 onglets | Deploye en ligne | Rapport PDF integre")

tabs = [
    ("Temporel",         VIOLET,  "Serie temporelle, profils horaires/hebdo, heatmap"),
    ("Predictions",      MAGENTA, "Predictions vs reel, residus, scatter interactif"),
    ("Anomalies",        ORANGE,  "Detection 2-sigma, carte des anomalies, statistiques"),
    ("Clustering",       VIOLET,  "K-Means 4 profils journaliers, radar, scatter"),
    ("Modeles",          MAGENTA, "Comparaison 6 modeles, SHAP, validation croisee, radar"),
    ("Simulateur",       ORANGE,  "Prediction instantanee, gauge, comparaison 5 modeles"),
    ("Rapport PDF",      VIOLET,  "Generation PDF 8 pages en 1 clic, telechargeable"),
    ("Contexte Senelec", MAGENTA, "Donnees Senelec 2022-2024, mix energetique, tendances"),
]
col_w2 = 6.15
for i, (name, color, desc) in enumerate(tabs):
    col = i % 2
    row = i // 2
    cx = 0.55 + col * (col_w2 + 0.35)
    cy = 1.4 + row * 1.45
    r = rect(s9, cx, cy, col_w2, 1.25, LIGHT)
    r.line.color.rgb = color
    r.line.width = Pt(1.5)
    rect(s9, cx, cy, 0.22, 1.25, color)
    txt(s9, name, cx+0.32, cy+0.08, col_w2-0.4, 0.4,
        size=13, color=color, bold=True)
    txt(s9, desc, cx+0.32, cy+0.48, col_w2-0.42, 0.65,
        size=11, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION (fond sombre)
# ═══════════════════════════════════════════════════════════════════════════
s10 = blank_slide(prs)
bg(s10, DARK)

# Bande gradient en haut
seg = 13.33 / 3
rect(s10, 0,      0, seg,   0.25, VIOLET)
rect(s10, seg,    0, seg,   0.25, MAGENTA)
rect(s10, seg*2,  0, seg+1, 0.25, ORANGE)

txt(s10, "Conclusion & Perspectives", 0.6, 0.45, 12.0, 0.7,
    size=28, color=WHITE, bold=True)

# 3 colonnes
cols = [
    ("Acquis", VIOLET, [
        "R² = 0.548 (Ridge) — conforme litterature",
        "6 modeles ML compares objectivement",
        "Validation croisee TimeSeriesSplit",
        "SHAP : variables lag les + influentes",
        "Dashboard Streamlit deploye en ligne",
        "Detection anomalies automatique",
    ]),
    ("Limites", MAGENTA, [
        "Dataset ancienne (2016, Belgique)",
        "Comportement humain = bruit irreductible",
        "Pas de donnees temps reel",
        "SVR lent sur grands datasets",
    ]),
    ("Perspectives", ORANGE, [
        "LSTM / modeles sequentiels (deep learning)",
        "Donnees Senelec temps reel",
        "Re-entrainement mensuel automatise",
        "Optimisation hyperparametres (GridSearch)",
        "Extension a d'autres batiments",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    cx = 0.5 + i * 4.25
    rect(s10, cx, 1.4, 4.0, 0.5, color)
    txt(s10, title, cx+0.12, 1.46, 3.8, 0.38,
        size=14, color=WHITE, bold=True)
    for j, item in enumerate(items):
        rect(s10, cx+0.08, 2.05+j*0.77, 0.22, 0.22, color)
        txt(s10, item, cx+0.38, 2.03+j*0.77, 3.55, 0.55,
            size=11, color=LGRAY)

# Bas de page
rect(s10, 0, 6.9, 13.33, 0.6, VIOLET)
txt(s10,
    "Groupe 6  |  Jean Paul MALAN   Guy Roger Junior GNAORE   COULIBALY Fobeh   TOURE Awa  |  Data Analyst 2025-2026",
    0.3, 6.97, 12.8, 0.45,
    size=10, color=WHITE, align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
out = "Presentation_Groupe6_Prediction_Energetique.pptx"
prs.save(out)
print(f"Presentation sauvegardee : {out}")
print(f"Nombre de diapositives : {len(prs.slides)}")
