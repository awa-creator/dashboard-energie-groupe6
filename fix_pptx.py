# -*- coding: utf-8 -*-
"""Corrections ciblees du PowerPoint par index de shape."""
from pptx import Presentation

prs = Presentation("Presentation_Groupe6_Prediction_Energetique.pptx")

def set_shape_text(slide, shape_idx, new_text):
    """Remplace le texte d'un shape en preservant le style du premier run."""
    shape = slide.shapes[shape_idx]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Vider tous les paragraphes sauf le premier
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # Mettre le nouveau texte dans le premier run du premier paragraphe
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
    elif tf.paragraphs:
        from pptx.util import Pt
        run = tf.paragraphs[0].add_run()
        run.text = new_text

def clear_shape(slide, shape_idx):
    set_shape_text(slide, shape_idx, "")

# ── SLIDE 4 : Supprimer Lasso et SVR ────────────────────────────────────
s4 = prs.slides[3]
clear_shape(s4, 33)  # "Lasso (L1) — alpha=0.1"
clear_shape(s4, 35)  # "SVR — RBF, C=100"
print("Slide 4 : Lasso et SVR supprimes")

# ── SLIDE 5 : Corriger titre, unites, texte residuel ────────────────────
s5 = prs.slides[4]
set_shape_text(s5, 4, "Resultats — Comparaison des 4 modeles")
set_shape_text(s5, 5, "Jeu de test chronologique (20% | 3 098 obs. horaires | unite : kWh)")
# Corriger "RMSE = 0.2970 Wh" -> kWh
set_shape_text(s5, 82, "RMSE = 0.2970 kWh   |   MAE = 0.1742 kWh")
# Corriger le paragraphe residuel "humain — bruit..."
shape80 = s5.shapes[80]
if shape80.has_text_frame:
    paras = shape80.text_frame.paragraphs
    # para[3] = "humain — bruit irreductible par definition."
    if len(paras) > 3:
        for run in paras[3].runs:
            run.text = ""
print("Slide 5 : titre, unites et texte corriges")

# ── SLIDE 6 : Mettre a jour les valeurs CV ───────────────────────────────
s6 = prs.slides[5]
# Les valeurs CV avec le caractere special —
set_shape_text(s6, 35, "0.4063  +/- 0.038")   # Ridge
set_shape_text(s6, 38, "0.3456  +/- 0.042")   # Gradient Boosting
set_shape_text(s6, 41, "0.4083  +/- 0.035")   # Random Forest
set_shape_text(s6, 44, "0.3929  +/- 0.040")   # Reg. Lineaire
# Supprimer les lignes SVR et doublon Reg. Lineaire
clear_shape(s6, 47)   # "0.465 +/- 0.044" (SVR)
clear_shape(s6, 49)   # "Reg. Lineaire" (doublon)
clear_shape(s6, 50)   # "0.451 +/- 0.040" (doublon)
print("Slide 6 : valeurs CV mises a jour")

# ── SLIDE 10 : Corriger R² = 0.548 ──────────────────────────────────────
s10 = prs.slides[9]
set_shape_text(s10, 7, "R2 = 0.4527 (Ridge) — dataset multi-batiments")
print("Slide 10 : R2 corrige")

prs.save("Presentation_Groupe6_Prediction_Energetique.pptx")
print("\nPresentation sauvegardee avec succes !")
